from __future__ import annotations

from pathlib import Path
from uuid import uuid4

from pptx import Presentation

from domain.document import (
    Document,
    DocumentMetadata,
    DocumentType,
    Page,
)
from infrastructure.parsers.base_parser import BaseParser


class PPTXParser(BaseParser):
    """
    Parser for Microsoft PowerPoint (.pptx) presentations.
    """

    @property
    def supported_extensions(self) -> tuple[str, ...]:
        return (".pptx",)

    def _extract(self, file_path: Path) -> Document:

        presentation = Presentation(file_path)

        pages: list[Page] = []

        all_text: list[str] = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):

            slide_text: list[str] = []

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text = shape.text.strip()

                    if text:

                        slide_text.append(text)

            combined = "\n".join(slide_text)

            pages.append(
                Page(
                    number=slide_number,
                    text=combined,
                )
            )

            all_text.append(combined)

        properties = presentation.core_properties

        return Document(
            id=str(uuid4()),
            filename=file_path.name,
            filepath=file_path,
            document_type=DocumentType.PPTX,
            pages=pages,
            metadata=DocumentMetadata(
                title=properties.title or "",
                author=properties.author or "",
                subject=properties.subject or "",
                keywords=[],
                language="",
            ),
            extracted_text="\n\n".join(all_text),
            page_count=len(pages),
            file_size=file_path.stat().st_size,
        )